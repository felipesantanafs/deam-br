"""
Gera deam24h_slides_editavel.pptx com elementos nativos do PowerPoint.
Todos os textos, formas e cores são editáveis diretamente no PowerPoint.

Requer: pip install python-pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ───────────────────────────────────────────────────────────────
ROOT  = Path(__file__).resolve().parent
FIGS  = ROOT.parent / "figuras_causal"
OUT   = ROOT / "deam24h_slides_editavel.pptx"

# ── Palette ─────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0E, 0x1B, 0x33)
NAVY2   = RGBColor(0x16, 0x29, 0x4C)
CREAM   = RGBColor(0xF4, 0xF1, 0xE8)
CREAMTX = RGBColor(0xEC, 0xE7, 0xD8)
GOLD    = RGBColor(0xC1, 0x9A, 0x3E)
INK     = RGBColor(0x2C, 0x2C, 0x2C)
SLATE   = RGBColor(0x5C, 0x6B, 0x82)
CARD    = RGBColor(0xFB, 0xFA, 0xF5)
HAIR    = RGBColor(0xD8, 0xD2, 0xC2)
ARED    = RGBColor(0xA8, 0x44, 0x3C)
AGREEN  = RGBColor(0x4F, 0x7A, 0x4E)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── Presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════

def _inch(x): return Inches(x)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=None, lc=None, lw=Pt(0.5)):
    """Add rectangle. fill/lc = RGBColor or None."""
    shp = slide.shapes.add_shape(1, _inch(x), _inch(y), _inch(w), _inch(h))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if lc:
        shp.line.color.rgb = lc
        shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, color, width_pt=0.5):
    """Add a straight connector line."""
    conn = slide.shapes.add_connector(
        1,                          # MSO_CONNECTOR.STRAIGHT
        _inch(x1), _inch(y1),
        _inch(x2), _inch(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width_pt)
    return conn


def txbox(slide, x, y, w, h):
    """Return (shape, text_frame) for a new textbox."""
    shp = slide.shapes.add_textbox(_inch(x), _inch(y), _inch(w), _inch(h))
    shp.word_wrap = True
    tf = shp.text_frame
    tf.word_wrap = True
    return shp, tf


def add_text(slide, text, x, y, w, h,
             font="Calibri", size=10, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT):
    """Add single-run textbox."""
    _, tf = txbox(slide, x, y, w, h)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _fmt(run, font, size, bold, italic, color)
    return tf


def add_para(tf, text, font="Calibri", size=10, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, new_para=False):
    """Append a paragraph (or use existing first) to a text_frame."""
    if new_para:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _fmt(run, font, size, bold, italic, color)
    return p


def _fmt(run, font, size, bold, italic, color):
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(str(path), _inch(x), _inch(y), _inch(w), _inch(h))


# ── Composite helpers ────────────────────────────────────────────────────

def add_footer(slide, num, total=6):
    add_text(slide, "•", 0.58, 7.17, 0.20, 0.26,
             font="Calibri", size=9, color=GOLD)
    add_text(slide, "AVALIAÇÃO DE IMPACTO  ·  DEAMS 24H NO BRASIL",
             0.80, 7.17, 9.5, 0.26,
             font="Calibri", size=7, bold=True, color=SLATE)
    pg = f"0{num} / 0{total}" if num < 10 else f"{num} / {total}"
    add_text(slide, pg, 11.5, 7.17, 1.6, 0.26,
             font="Georgia", size=9, italic=True, color=SLATE,
             align=PP_ALIGN.RIGHT)


def add_header(slide, kicker, title, subtitle):
    """Standard header: kicker line + hairline + title box with gold left border."""
    # Kicker row
    add_text(slide, "•", 0.55, 0.31, 0.20, 0.22,
             font="Calibri", size=8, bold=True, color=GOLD)
    add_text(slide, kicker.upper(), 0.75, 0.31, 7.5, 0.22,
             font="Calibri", size=7, bold=True, color=NAVY2)
    add_text(slide, "DEAM · BR", 9.3, 0.31, 3.7, 0.22,
             font="Georgia", size=8, italic=True, color=SLATE,
             align=PP_ALIGN.RIGHT)
    # Hairline
    add_rect(slide, 0.55, 0.57, 12.22, 0.012, fill=HAIR)
    # Gold left border strip for title
    add_rect(slide, 0.55, 0.63, 0.042, 0.55, fill=GOLD)
    # Title
    add_text(slide, title, 0.67, 0.63, 12.0, 0.33,
             font="Georgia", size=17, bold=True, color=NAVY2)
    # Subtitle
    add_text(slide, subtitle, 0.67, 0.97, 12.0, 0.22,
             font="Georgia", size=9.5, italic=True, color=SLATE)


def add_eyebrow(slide, text, x, y):
    add_text(slide, text.upper(), x, y, 6.5, 0.22,
             font="Calibri", size=7, bold=True, color=GOLD)


def add_callout(slide, title, body, x, y, w, h, bc=GOLD):
    """Colored left-border callout card."""
    add_rect(slide, x, y, w, h, fill=CARD)
    add_rect(slide, x, y, 0.042, h, fill=bc)       # left border
    add_text(slide, title.upper(), x+0.11, y+0.055, w-0.15, 0.20,
             font="Calibri", size=7, bold=True, color=bc)
    add_text(slide, body, x+0.11, y+0.25, w-0.15, h-0.28,
             font="Calibri", size=9, color=INK)


def add_chainstep(slide, title, subtitle, x, y, w, h, bc=NAVY2):
    """Causal chain box with colored left border."""
    add_rect(slide, x, y, w, h, fill=CARD, lc=HAIR, lw=Pt(0.4))
    add_rect(slide, x, y, 0.038, h, fill=bc)
    title_y = y + (h - 0.22) / 2 if not subtitle else y + 0.05
    add_text(slide, title, x+0.10, title_y, w-0.14, 0.24,
             font="Georgia", size=11, bold=True, color=NAVY2,
             align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, subtitle, x+0.10, y+0.29, w-0.14, 0.20,
                 font="Calibri", size=8, italic=True, color=SLATE,
                 align=PP_ALIGN.CENTER)


def add_readitem(slide, term, desc, x, y, w):
    """Gold bullet + bold term + normal description (multi-run)."""
    add_text(slide, "•", x, y, 0.20, 0.26, font="Calibri", size=10, color=GOLD)
    _, tf = txbox(slide, x+0.19, y, w-0.19, 0.28)
    p = tf.paragraphs[0]
    r1 = p.add_run(); _fmt(r1, "Calibri", 9, True,  False, NAVY2); r1.text = term + "  "
    r2 = p.add_run(); _fmt(r2, "Calibri", 9, False, False, INK);   r2.text = desc


def add_bullet(slide, bold_part, rest, x, y, w):
    """Gold dash bullet + bold + normal body."""
    # Gold dash (em dash as bullet)
    add_text(slide, "—", x, y, 0.22, 0.28, font="Calibri", size=9, color=GOLD)
    _, tf = txbox(slide, x+0.22, y, w-0.22, 0.34)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if bold_part:
        r1 = p.add_run(); _fmt(r1, "Calibri", 9, True, False, NAVY2); r1.text = bold_part + " "
    if rest:
        r2 = p.add_run(); _fmt(r2, "Calibri", 9, False, False, INK); r2.text = rest


def add_chartcard(slide, img_path, x, y, w, h):
    """White card with hairline border + image."""
    add_rect(slide, x, y, w, h, fill=WHITE, lc=HAIR, lw=Pt(0.5))
    pad = 0.07
    add_image(slide, img_path, x+pad, y+pad, w-2*pad, h-2*pad)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 0 — CAPA (NAVY)
# ══════════════════════════════════════════════════════════════════════════
s0 = prs.slides.add_slide(BLANK)
set_bg(s0, NAVY)

# Filetes verticais dourados
add_line(s0, 13.333-1.70, 0.0, 13.333-1.70, 7.5, GOLD, 0.5)
add_line(s0, 13.333-1.30, 0.0, 13.333-1.30, 7.5, GOLD, 0.5)

# Top: bullet + "DEAM · BR"
add_text(s0, "•", 0.60, 0.32, 0.22, 0.26, font="Calibri", size=9, color=GOLD)
add_text(s0, "DEAM · BR", 0.84, 0.32, 5.0, 0.26,
         font="Georgia", size=9, italic=True, color=CREAMTX)

# Eyebrow
add_text(s0, "AVALIAÇÃO DE IMPACTO CAUSAL",
         0.60, 2.55, 11.0, 0.38,
         font="Calibri", size=10.5, bold=True, color=GOLD)

# Main title — two paragraphs
_, tf_t = txbox(s0, 0.60, 2.98, 10.5, 1.85)
p1 = tf_t.paragraphs[0]
r  = p1.add_run(); _fmt(r, "Georgia", 34, True, False, CREAMTX)
r.text = "Conversão de DEAMs para"
p2 = tf_t.add_paragraph()
r2 = p2.add_run(); _fmt(r2, "Georgia", 34, True, False, CREAMTX)
r2.text = "Plantão 24h no Brasil"

# Gold rule
add_rect(s0, 0.60, 4.88, 2.2, 0.038, fill=GOLD)

# Subtitle
add_text(s0, "Do acesso à proteção — uma abordagem via Callaway & Sant'Anna (2021)",
         0.60, 4.96, 10.5, 0.36,
         font="Georgia", size=11, italic=True, color=CREAMTX)

# Authors (multi-run)
_, tf_a = txbox(s0, 0.60, 5.40, 10.0, 0.34)
pa = tf_a.paragraphs[0]
r_a1 = pa.add_run(); _fmt(r_a1, "Calibri", 13, False, False, CREAMTX); r_a1.text = "Felipe Silva  "
r_a2 = pa.add_run(); _fmt(r_a2, "Calibri", 13, False, False, GOLD);    r_a2.text = "·"
r_a3 = pa.add_run(); _fmt(r_a3, "Calibri", 13, False, False, CREAMTX); r_a3.text = "  Rafael Issa de Lima"

# Bottom label
add_text(s0, "APRESENTAÇÃO TÉCNICA  ·  DEAM-BR  ·  PAINEL 2009–2019",
         0.60, 7.04, 11.0, 0.30,
         font="Calibri", size=7, bold=True, color=GOLD)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — O PROBLEMA E A CADEIA CAUSAL
# ══════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
set_bg(s1, CREAM)
add_header(s1, "Contexto · 01",
           "O Problema e a Cadeia Causal",
           "Por que ampliar o horário de atendimento das delegacias?")
add_footer(s1, 1)

# LEFT (x=0.55, w=5.85)
add_eyebrow(s1, "O Problema", 0.55, 1.30)
add_text(s1, "DEAMs operam, em regra, em horário comercial (9h–18h). Mas a violência não:",
         0.55, 1.56, 5.75, 0.44,
         font="Calibri", size=9.5, color=INK)

# 57% hero
add_text(s1, "57%", 0.55, 2.08, 5.75, 1.55,
         font="Georgia", size=76, bold=True, color=ARED)
add_text(s1, "das notificações ocorrem fora do horário comercial (08h–17h)",
         0.55, 3.70, 5.30, 0.44,
         font="Calibri", size=9, italic=True, color=SLATE)

# RIGHT — causal chain (x=6.68, w=6.32)
CX, CW = 6.68, 6.32
add_eyebrow(s1, "Cadeia causal (cifra oculta)", CX, 1.30)

steps = [
    ("Plantão 24h",    "",                              GOLD),
    ("↑ Acesso",       "mais notificações noturnas",    AGREEN),
    ("Proteção",       "medidas protetivas de urgência",AGREEN),
    ("↓ Letalidade",   "feminicídios",                  ARED),
]
cy = 1.55
for title, sub, bc in steps:
    sh = 0.52 if sub else 0.38
    add_chainstep(s1, title, sub, CX, cy, CW, sh, bc)
    cy += sh
    if title != "↓ Letalidade":
        add_text(s1, "↓", CX + CW/2 - 0.15, cy + 0.01, 0.30, 0.20,
                 font="Calibri", size=13, color=GOLD, align=PP_ALIGN.CENTER)
        cy += 0.22

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — METODOLOGIA
# ══════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
set_bg(s2, CREAM)
add_header(s2, "Metodologia · 02",
           "Metodologia e o Desafio dos Dados",
           "Identificação sob adoção escalonada — desfechos em taxa /100k hab.")
add_footer(s2, 2)

# LEFT
add_callout(s2,
    "O problema do TWFE",
    "Goodman-Bacon (2021): com adoção escalonada (10 coortes), o estimador "
    "de efeitos fixos gera pesos negativos sob efeitos heterogêneos.",
    0.55, 1.30, 5.85, 1.18, bc=ARED)

add_callout(s2,
    "A solução · CS-DiD (2021)",
    "Efeitos por coorte-tempo ATT(g,t): os 38 tratados comparados apenas "
    "com os 247 nunca-tratados (grupo de controle limpo).",
    0.55, 2.58, 5.85, 1.08, bc=AGREEN)

# RIGHT
add_eyebrow(s2, "38 tratados em 10 coortes (2014 = maior)", 6.70, 1.30)
add_chartcard(s2, FIGS / "00b_coortes.png", 6.70, 1.56, 6.30, 4.52)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESULTADOS I: ACESSO
# ══════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
set_bg(s3, CREAM)
add_header(s3, "Resultados · 03",
           "Acesso e o Viés de Tendência",
           "Desfecho: taxa de notificações /100k (SINAN)")
add_footer(s3, 3)

add_chartcard(s3, FIGS / "notificacoes_02_event_study.png",
              0.55, 1.30, 7.52, 4.78)

# Right panel
add_eyebrow(s3, "Leitura científica", 8.35, 1.30)
ry = 1.56
add_readitem(s3, "ATT = −20,51",           "(p = 0,001)",           8.35, ry, 4.72); ry += 0.36
add_readitem(s3, "Pré-tendências violadas", "p Wald = 0,004",        8.35, ry, 4.72); ry += 0.36
add_readitem(s3, "Pontos em t<0",           "já divergem de zero",   8.35, ry, 4.72); ry += 0.50

add_callout(s3, "Veredicto",
    "O efeito capta uma trajetória prévia, não a política. "
    "Não interpretável como causal.",
    8.35, ry, 4.72, 1.02, bc=ARED)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RESULTADOS II: LETALIDADE
# ══════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
set_bg(s4, CREAM)
add_header(s4, "Resultados · 04",
           "Letalidade e Adoção Reativa",
           "Desfecho: taxa de feminicídios /100k (SIM)")
add_footer(s4, 4)

add_chartcard(s4, FIGS / "feminicidios_02_event_study.png",
              0.55, 1.30, 7.52, 4.78)

add_eyebrow(s4, "Leitura científica", 8.35, 1.30)
ry = 1.56
add_readitem(s4, "Pré-tendências OK",  "p Wald = 0,165",              8.35, ry, 4.72); ry += 0.38
add_readitem(s4, "ATT = +0,438",       "(p = 0,082) — sinal invertido", 8.35, ry, 4.72); ry += 0.55

add_callout(s4, "Intuição · Adoção Reativa",
    "Gestores abrem DEAMs 24h onde a letalidade já está em viés de alta "
    "— causalidade reversa (endogeneidade).",
    8.35, ry, 4.72, 1.14, bc=GOLD)

# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MECANISMO & ROADMAP
# ══════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
set_bg(s5, CREAM)
add_header(s5, "Mecanismo & Agenda · 05",
           "Mecanismo e o Caminho à Frente",
           "Evidência de sazonalidade e próximos passos metodológicos")
add_footer(s5, 5)

# LEFT — table
add_eyebrow(s5, "A demanda noturna é constante", 0.55, 1.30)

# Table header
add_rect(s5, 0.55, 1.54, 5.95, 0.32, fill=NAVY2)
add_text(s5, "Grupo / Período",    0.62, 1.56, 3.8, 0.28,
         font="Calibri", size=8.5, bold=True, color=CREAM)
add_text(s5, "% fora 08–17h",      4.48, 1.56, 1.9, 0.28,
         font="Calibri", size=8.5, bold=True, color=CREAM, align=PP_ALIGN.RIGHT)

# Table rows
trows = [
    ("Tratados — antes do 24h", "57,1%", CARD),
    ("Tratados — após o 24h",   "55,9%", WHITE),
    ("Controles (comercial)",   "57,6%", CARD),
]
ry = 1.86
for label, val, bg in trows:
    add_rect(s5, 0.55, ry, 5.95, 0.32, fill=bg, lc=HAIR, lw=Pt(0.3))
    add_text(s5, label, 0.62, ry+0.05, 3.80, 0.24,
             font="Calibri", size=9, color=INK)
    add_text(s5, val,   4.48, ry+0.05, 1.90, 0.24,
             font="Calibri", size=9, bold=True, color=ARED, align=PP_ALIGN.RIGHT)
    ry += 0.32

# italic note below table
add_text(s5, "~56% fora do horário comercial antes e depois do 24h.",
         0.55, ry+0.14, 5.95, 0.30,
         font="Georgia", size=9, italic=True, color=SLATE)

# RIGHT — roadmap + callout
add_eyebrow(s5, "Roadmap · Próximos Passos", 6.85, 1.30)
by = 1.56
bullets = [
    ("DR-DiD",                     "— estimador duplamente robusto."),
    ("Covariáveis socioeconômicas", "(IDH, renda, urbanização) → tendências paralelas condicionais."),
    ("Controle de antecipação",     "e grupo de não-ainda-tratados."),
]
for bp, rest in bullets:
    add_bullet(s5, bp, rest, 6.85, by, 6.05)
    by += 0.44

add_callout(s5, "Mensagem Central",
    "Sem ajustar para adoção reativa e tendências prévias, avaliações "
    "de plantão 24h tendem a ser pessimistas por construção.",
    6.85, by+0.18, 6.05, 1.05, bc=NAVY2)

# ══════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"OK -> {OUT}  (6 slides nativos editaveis)")
