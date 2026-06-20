"""Monta deam24h_slides.pptx a partir dos PNGs renderizados do PDF Beamer.
Cada slide do PowerPoint = uma página do PDF, em full-bleed 16:9 (idêntico ao PDF)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
RENDER = ROOT / "render"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]       # layout em branco

pngs = sorted(RENDER.glob("slide-*.png"), key=lambda p: int(p.stem.split("-")[1]))
for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0,
                             width=prs.slide_width, height=prs.slide_height)

out = ROOT / "deam24h_slides.pptx"
prs.save(out)
print(f"OK -> {out}  ({len(pngs)} slides)")
